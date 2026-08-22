from pptx import Presentation


def extract_text_from_pptx(file):
    presentation = Presentation(file)

    slides = []

    for slide_number, slide in enumerate(presentation.slides, start=1):

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_text.append(text)

        slides.append({
            "slide_number": slide_number,
            "text": "\n".join(slide_text)
        })

    return slides